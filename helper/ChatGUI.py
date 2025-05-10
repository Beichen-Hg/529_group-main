import os
import threading
import queue
import time
from datetime import datetime
from dotenv import load_dotenv
import speech_recognition as sr
import pyttsx3
import tkinter as tk
from tkinter import ttk, scrolledtext, messagebox, filedialog
from tkinterdnd2 import DND_FILES, TkinterDnD
from langchain_openai import ChatOpenAI
from langchain.memory import ConversationBufferMemory
from langchain.chains import ConversationChain
import fitz  # PyMuPDF for PDF
from docx import Document  # python-docx for Word
from pptx import Presentation  # python-pptx for PPT
import base64
from PIL import Image
from io import BytesIO
import cv2
from PIL import ImageTk
import numpy as np
import matplotlib.pyplot as plt
from matplotlib.backends.backend_tkagg import FigureCanvasTkAgg
import seaborn as sns
from matplotlib.backends.backend_tkagg import NavigationToolbar2Tk

class ChatGUI(TkinterDnD.Tk):
    def __init__(self, bot):
        # 初始化队列
        self.queue = queue.Queue()
        self.bot = bot
        self.pending_files = []  # 存储待处理的文件信息
        self.logged_in_user = None # To store the username of the logged-in user
        
        # 添加新的实例变量
        self.preview_window = None
        self.preview_canvas = None
        self.preview_fig = None
        self.recording_preview = False
        
        # 调用父类构造函数
        super().__init__()
        
        # 设置主题色
        self.style = ttk.Style()
        
        # 配置基础样式
        self.style.configure('Modern.TFrame', background='#f0f0f0')
        self.style.configure('Modern.TLabel', 
                           background='#f0f0f0', 
                           font=('Microsoft YaHei', 10))
        self.style.configure('Modern.TButton', 
                           font=('Microsoft YaHei', 10),
                           padding=5)
        self.style.configure('Modern.TEntry',
                           font=('Microsoft YaHei', 10),
                           padding=5)
        
        # 配置LabelFrame样式
        self.style.configure('Modern.TLabelframe', 
                           background='#f0f0f0',
                           font=('Microsoft YaHei', 10))
        self.style.configure('Modern.TLabelframe.Label', 
                           background='#f0f0f0',
                           font=('Microsoft YaHei', 10))
        
        self.title("阿kie")
        self._setup_ui()
        self.after(100, self.check_queue)
        
        # 启动持续语音监听
        self.bot.start_listening(self._handle_voice_input)
        
        # 设置文件拖放
        self.drop_target_register(DND_FILES)
        self.dnd_bind('<<Drop>>', self._handle_drop)

    def check_queue(self):
        """检查并处理消息队列"""
        try:
            while True:
                func = self.queue.get_nowait()
                func()
        except queue.Empty:
            pass
        finally:
            self.after(100, self.check_queue)
            
    def _handle_voice_input(self, text):
        """处理语音输入"""
        print(f"收到语音输入: {text}")
        
        # 如果正在处理或说话，跳过这次输入
        if self.bot.is_processing or self.bot.is_speaking:
            print("系统正在处理或说话，跳过此次输入")
            return
            
        # 显示语音输入
        self._thread_safe_display("您", text)
        
        # 处理语音输入
        threading.Thread(
            target=self._process_query,
            args=(text,),
            daemon=True
        ).start()
        
    def _setup_ui(self):
        """构建用户界面"""
        # 设置窗口最小尺寸
        self.minsize(800, 650) # Slightly increase min height for user controls
        
        # 创建主框架
        main_frame = ttk.Frame(self, style='Modern.TFrame')
        main_frame.pack(fill=tk.BOTH, expand=True, padx=10, pady=10)
        
        # 创建左右分隔框架
        left_frame = ttk.Frame(main_frame, style='Modern.TFrame')
        left_frame.pack(side=tk.LEFT, fill=tk.BOTH, expand=True, padx=(0, 5))
        left_frame.pack_propagate(False)
        left_frame.configure(width=300)  # 设置左侧宽度为总宽度的30%
        
        right_frame = ttk.Frame(main_frame, style='Modern.TFrame')
        right_frame.pack(side=tk.LEFT, fill=tk.BOTH, expand=True, padx=(5, 0))
        
        # 左侧区域：摄像头和文件上传
        self._build_left_section(left_frame)
        
        # 右侧区域：聊天记录和输入控制
        self._build_right_section(right_frame)

    def _build_left_section(self, parent):
        """构建左侧区域（摄像头和文件上传）"""
        # 创建上下分隔框架
        top_frame = ttk.Frame(parent, style='Modern.TFrame')
        top_frame.pack(side=tk.TOP, fill=tk.BOTH, expand=True, pady=(0, 5))
        
        bottom_frame = ttk.Frame(parent, style='Modern.TFrame')
        bottom_frame.pack(side=tk.TOP, fill=tk.BOTH, expand=True, pady=(5, 0))
        
        # 摄像头区域
        camera_frame = ttk.LabelFrame(top_frame, text="摄像头", style='Modern.TLabelframe')
        camera_frame.pack(fill=tk.BOTH, expand=True, padx=5, pady=5)
        
        # 创建一个固定大小的Frame来容纳摄像头显示
        camera_container = ttk.Frame(camera_frame, style='Modern.TFrame')
        camera_container.pack(padx=5, pady=5)
        camera_container.pack_propagate(False)
        camera_container.configure(width=300, height=225)
        
        # 摄像头画面显示区域
        self.camera_display = ttk.Label(camera_container, style='Modern.TLabel')
        self.camera_display.pack(fill=tk.BOTH, expand=True)
        
        # 进度条
        self.progress_frame = ttk.Frame(camera_frame, style='Modern.TFrame')
        self.progress_frame.pack(fill=tk.X, padx=5, pady=(0, 5))
        self.progress_bar = ttk.Progressbar(
            self.progress_frame,
            mode='determinate',
            length=300
        )
        self.progress_bar.pack(fill=tk.X)
        self.progress_bar.pack_forget()  # 初始时隐藏
        
        # 摄像头控制按钮框架
        camera_control_frame = ttk.Frame(camera_frame, style='Modern.TFrame')
        camera_control_frame.pack(fill=tk.X, pady=5)
        
        # 按钮框架
        button_frame = ttk.Frame(camera_control_frame, style='Modern.TFrame')
        button_frame.pack(fill=tk.X)
        
        # 摄像头控制按钮
        self.camera_button = ttk.Button(
            button_frame,
            text="开启摄像头",
            command=self._toggle_camera,
            style='Modern.TButton'
        )
        self.camera_button.pack(side=tk.LEFT, padx=5)
        
        # 语音识别控制按钮 (修改为长按)
        self.voice_button = ttk.Button(
            button_frame,
            text="按住说话",
            style='Modern.TButton'
        )
        self.voice_button.pack(side=tk.LEFT, padx=5)
        self.voice_button.bind("<ButtonPress-1>", self._on_voice_button_press)
        self.voice_button.bind("<ButtonRelease-1>", self._on_voice_button_release)
        
        # 性格分析按钮
        self.personality_button = ttk.Button(
            button_frame,
            text="性格分析",
            command=self._start_personality_analysis,
            style='Modern.TButton'
        )
        self.personality_button.pack(side=tk.LEFT, padx=5)
        
        # 历史记录按钮
        self.history_button = ttk.Button(
            button_frame,
            text="历史记录",
            command=self._show_history,
            style='Modern.TButton'
        )
        self.history_button.pack(side=tk.LEFT, padx=5)
        
        # 文件上传区域
        upload_frame = ttk.LabelFrame(bottom_frame, text="文件上传", style='Modern.TLabelframe')
        upload_frame.pack(fill=tk.BOTH, expand=True, padx=5, pady=5)
        
        # 创建容器框架用于居中显示
        center_frame = ttk.Frame(upload_frame, style='Modern.TFrame')
        center_frame.pack(expand=True)
        
        # 文件上传按钮
        self.upload_button = ttk.Button(
            center_frame,
            text="选择文件",
            command=self._open_file_dialog,
            style='Modern.TButton'
        )
        self.upload_button.pack(pady=(0, 5))
        
        # 添加拖放提示标签
        self.drop_label = ttk.Label(
            center_frame,
            text="或将文件拖放到此处",
            wraplength=250,
            style='Modern.TLabel'
        )
        self.drop_label.pack()
        
        # 文件处理状态显示区域
        self.file_status_frame = ttk.Frame(upload_frame, style='Modern.TFrame')
        self.file_status_frame.pack(fill=tk.X, padx=5, pady=5)
        
        # 文件处理状态标签
        self.file_status_label = ttk.Label(
            self.file_status_frame,
            text="",
            wraplength=250,
            style='Modern.TLabel'
        )
        self.file_status_label.pack()

    def _build_right_section(self, parent):
        """构建右侧区域（聊天记录和输入控制）"""
        # 新增：用户认证控制区域
        self._build_user_controls(parent)

        # 聊天记录区
        self._build_chat_area(parent)
        
        # 输入控制区
        self._build_input_controls(parent)

        # 初始化UI的登录状态
        self._update_ui_for_login_status()

    def _build_chat_area(self, parent):
        """构建聊天记录显示区域"""
        chat_frame = ttk.Frame(parent, style='Modern.TFrame')
        chat_frame.pack(expand=True, fill=tk.BOTH, pady=(0,5)) 
        
        self.chat_area = scrolledtext.ScrolledText(
            chat_frame,
            wrap=tk.WORD,
            state='disabled',
            font=('Microsoft YaHei', 10),
            spacing3=5,
            bg='#ffffff',
            fg='#333333',
            padx=10,
            pady=10
        )
        self.chat_area.pack(expand=True, fill=tk.BOTH, padx=5, pady=0) 
        
        # Initial message can be set here or after login
        # self._thread_safe_display("系统", "请登录或注册以开始对话。")

    def _build_input_controls(self, parent):
        """构建输入控制区域"""
        control_frame = ttk.Frame(parent, style='Modern.TFrame')
        control_frame.pack(fill=tk.X, pady=(0,5), padx=5) 
        
        self.input_text = ttk.Entry(control_frame, style='Modern.TEntry', font=('Microsoft YaHei', 10))
        self.input_text.pack(side=tk.LEFT, expand=True, fill=tk.X, padx=(0, 5))
        
        self.send_button = ttk.Button(
            control_frame,
            text="发送",
            command=self._on_text_input,
            style='Modern.TButton'
        )
        self.send_button.pack(side=tk.LEFT, padx=5)
        
        self.tts_button = ttk.Button(
            control_frame,
            text="语音播报: 开", 
            command=self._toggle_tts,
            style='Modern.TButton'
        )
        self.tts_button.pack(side=tk.LEFT, padx=5)
        
        self.input_text.bind('<Return>', self._on_text_input) 

    def _build_user_controls(self, parent):
        """构建用户登录/注册控制区域"""
        user_control_frame = ttk.Frame(parent, style='Modern.TFrame')
        user_control_frame.pack(fill=tk.X, pady=(5, 5), padx=5) # Pack above chat area

        self.user_status_label = ttk.Label(user_control_frame, text="未登录", style='Modern.TLabel', anchor="w")
        self.user_status_label.pack(side=tk.LEFT, padx=(0,10), fill=tk.X, expand=True)

        self.logout_button = ttk.Button(user_control_frame, text="登出", command=self._handle_logout, style='Modern.TButton')
        self.logout_button.pack(side=tk.RIGHT, padx=5)
        
        self.register_button = ttk.Button(user_control_frame, text="注册", command=self._show_register_dialog, style='Modern.TButton')
        self.register_button.pack(side=tk.RIGHT, padx=5)

        self.login_button = ttk.Button(user_control_frame, text="登录", command=self._show_login_dialog, style='Modern.TButton')
        self.login_button.pack(side=tk.RIGHT, padx=5)

    def _update_ui_for_login_status(self):
        """根据登录状态更新UI元素"""
        if self.logged_in_user:
            self.user_status_label.config(text=f"欢迎, {self.logged_in_user}!")
            if hasattr(self, 'login_button'): self.login_button.config(state='disabled')
            if hasattr(self, 'register_button'): self.register_button.config(state='disabled')
            if hasattr(self, 'logout_button'): self.logout_button.config(state='normal')
            # 聊天输入在登录后应可用
            if hasattr(self, 'input_text'): self.input_text.config(state='normal')
            if hasattr(self, 'send_button'): self.send_button.config(state='normal')
            if hasattr(self, 'voice_button'): self.voice_button.config(state='normal') # Enable voice button
        else:
            self.user_status_label.config(text="未登录")
            if hasattr(self, 'login_button'): self.login_button.config(state='normal')
            if hasattr(self, 'register_button'): self.register_button.config(state='normal')
            if hasattr(self, 'logout_button'): self.logout_button.config(state='disabled')
            # 聊天输入在未登录时禁用 (可选策略)
            # self.input_text.config(state='disabled')
            # self.send_button.config(state='disabled')
            # self.voice_button.config(state='disabled') # Disable voice button

    def _show_login_dialog(self):
        """显示登录对话框"""
        login_dialog = tk.Toplevel(self)
        login_dialog.title("用户登录")
        login_dialog.geometry("320x160") # Adjusted size
        login_dialog.resizable(False, False)
        login_dialog.transient(self) 
        login_dialog.grab_set() 

        dialog_frame = ttk.Frame(login_dialog, padding="10 10 10 10")
        dialog_frame.pack(expand=True, fill=tk.BOTH)

        ttk.Label(dialog_frame, text="用户名:").grid(row=0, column=0, padx=5, pady=5, sticky="w")
        username_entry = ttk.Entry(dialog_frame, width=25)
        username_entry.grid(row=0, column=1, padx=5, pady=5)

        ttk.Label(dialog_frame, text="密码:").grid(row=1, column=0, padx=5, pady=5, sticky="w")
        password_entry = ttk.Entry(dialog_frame, show="*", width=25)
        password_entry.grid(row=1, column=1, padx=5, pady=5)
        
        username_entry.focus_set()

        button_frame = ttk.Frame(dialog_frame)
        button_frame.grid(row=2, column=0, columnspan=2, pady=15)
        
        login_btn_dialog = ttk.Button(button_frame, text="登录", 
                   command=lambda: self._perform_login(username_entry.get(), password_entry.get(), login_dialog))
        login_btn_dialog.pack(side=tk.LEFT, padx=10)
        
        ttk.Button(button_frame, text="取消", command=login_dialog.destroy).pack(side=tk.LEFT, padx=10)
        
        login_dialog.bind('<Return>', lambda event: self._perform_login(username_entry.get(), password_entry.get(), login_dialog))
        # Center the dialog
        login_dialog.update_idletasks()
        x = self.winfo_x() + (self.winfo_width() // 2) - (login_dialog.winfo_width() // 2)
        y = self.winfo_y() + (self.winfo_height() // 2) - (login_dialog.winfo_height() // 2)
        login_dialog.geometry(f"+{x}+{y}")

    def _perform_login(self, username, password, dialog):
        """执行登录操作"""
        if not username or not password:
            messagebox.showerror("登录失败", "用户名和密码不能为空。", parent=dialog)
            return

        success, message = self.bot.login_user(username, password)
        if success:
            self.logged_in_user = username
            self._update_ui_for_login_status()
            # messagebox.showinfo("登录成功", message, parent=self) # Show on main window
            self._thread_safe_display("系统", f"用户 {username} 已登录。聊天记录已加载。")
            dialog.destroy()
        else:
            messagebox.showerror("登录失败", message, parent=dialog)
            # Keep dialog open for retry

    def _show_register_dialog(self):
        """显示注册对话框"""
        register_dialog = tk.Toplevel(self)
        register_dialog.title("用户注册")
        register_dialog.geometry("350x200") 
        register_dialog.resizable(False, False)
        register_dialog.transient(self)
        register_dialog.grab_set()

        dialog_frame = ttk.Frame(register_dialog, padding="10 10 10 10")
        dialog_frame.pack(expand=True, fill=tk.BOTH)

        ttk.Label(dialog_frame, text="用户名:").grid(row=0, column=0, padx=5, pady=5, sticky="w")
        reg_username_entry = ttk.Entry(dialog_frame, width=30)
        reg_username_entry.grid(row=0, column=1, padx=5, pady=5)

        ttk.Label(dialog_frame, text="密码:").grid(row=1, column=0, padx=5, pady=5, sticky="w")
        reg_password_entry = ttk.Entry(dialog_frame, show="*", width=30)
        reg_password_entry.grid(row=1, column=1, padx=5, pady=5)

        ttk.Label(dialog_frame, text="确认密码:").grid(row=2, column=0, padx=5, pady=5, sticky="w")
        reg_confirm_password_entry = ttk.Entry(dialog_frame, show="*", width=30)
        reg_confirm_password_entry.grid(row=2, column=1, padx=5, pady=5)
        
        reg_username_entry.focus_set()

        button_frame = ttk.Frame(dialog_frame)
        button_frame.grid(row=3, column=0, columnspan=2, pady=15)

        reg_btn_dialog = ttk.Button(button_frame, text="注册", 
                   command=lambda: self._perform_register(reg_username_entry.get(), 
                                                          reg_password_entry.get(), 
                                                          reg_confirm_password_entry.get(), 
                                                          register_dialog))
        reg_btn_dialog.pack(side=tk.LEFT, padx=10)
        ttk.Button(button_frame, text="取消", command=register_dialog.destroy).pack(side=tk.LEFT, padx=10)
        
        register_dialog.bind('<Return>', lambda event: self._perform_register(reg_username_entry.get(), 
                                                                          reg_password_entry.get(), 
                                                                          reg_confirm_password_entry.get(), 
                                                                          register_dialog))
        # Center the dialog
        register_dialog.update_idletasks()
        x = self.winfo_x() + (self.winfo_width() // 2) - (register_dialog.winfo_width() // 2)
        y = self.winfo_y() + (self.winfo_height() // 2) - (register_dialog.winfo_height() // 2)
        register_dialog.geometry(f"+{x}+{y}")

    def _perform_register(self, username, password, confirm_password, dialog):
        """执行注册操作"""
        if not username or not password or not confirm_password:
            messagebox.showerror("注册失败", "所有字段均不能为空。", parent=dialog)
            return
        if password != confirm_password:
            messagebox.showerror("注册失败", "两次输入的密码不匹配。", parent=dialog)
            return
        # Basic password complexity (example: at least 6 chars)
        if len(password) < 6:
            messagebox.showerror("注册失败", "密码长度至少需要6位。", parent=dialog)
            return

        success, message = self.bot.register_user(username, password)
        if success:
            messagebox.showinfo("注册成功", message + "\n现在您可以登录了。", parent=self) # Show on main
            dialog.destroy()
        else:
            messagebox.showerror("注册失败", message, parent=dialog)
            # Keep dialog open

    def _handle_logout(self):
        """处理登出操作"""
        if self.logged_in_user:
            user_to_logout = self.logged_in_user
            self.bot.logout_user() 
            self.logged_in_user = None
            self._update_ui_for_login_status()
            # messagebox.showinfo("已登出", f"用户 {user_to_logout} 已成功登出。\n聊天记录已保存。", parent=self)
            self._thread_safe_display("系统", f"用户 {user_to_logout} 已登出。聊天记录已保存。")
        else:
            messagebox.showwarning("登出", "当前没有用户登录。", parent=self)
            
    def _on_text_input(self, event=None):
        """处理文本输入"""
        # If a user is not logged in, optionally prevent chatting or use a guest session.
        # For now, we assume if not logged_in_user, bot.text_chat uses "default_session"
        # as per VoiceChatBot's current_user_id logic.
        # if not self.logged_in_user:
        #     messagebox.showwarning("请先登录", "请先登录后再开始聊天。", parent=self)
        #     return

        if self.bot.is_speaking: # This should call bot.stop_speaking() or similar
            # self.bot.stop_processing() # Ensure bot has a way to interrupt speech if needed
            if hasattr(self.bot, 'stop_speaking_immediately') and callable(self.bot.stop_speaking_immediately):
                self.bot.stop_speaking_immediately()
            elif hasattr(self.bot, 'speech_queue') and hasattr(self.bot.speech_queue, 'queue'): # Clear queue
                 with self.bot.speech_queue.mutex:
                    self.bot.speech_queue.queue.clear()

        text = self.input_text.get().strip()
        
        # user_id_to_use = self.bot.get_current_user_id() if self.bot.get_current_user_id() else "default_session"
        # VoiceChatBot's text_chat method already handles using self.current_user_id if set.

        if self.pending_files:
            original_prompt_text = text if text else "请分析这些文件的内容。"
            files_info = "\n".join([
                f"文件 {i+1}: {file['name']} ({file['type']})\n---BEGIN CONTENT---\n{file['content']}\n---END CONTENT---"
                for i, file in enumerate(self.pending_files)
            ])
            full_prompt = f"{original_prompt_text}\n\n附加文件内容如下：\n{files_info}"
            
            self.pending_files = []
            self._update_file_status("") # Clear file status
            
            self.queue.put(lambda: self.input_text.delete(0, tk.END))
            self._thread_safe_display("您", original_prompt_text) 
            
            threading.Thread(
                target=self._process_query,
                args=(full_prompt,),
                daemon=True
            ).start()
            return

        if not text:
            return
        
        self.queue.put(lambda: self.input_text.delete(0, tk.END))
        self._thread_safe_display("您", text)
        
        threading.Thread(
            target=self._process_query,
            args=(text,),
            daemon=True
        ).start()

    def _process_query(self, query):
        """处理用户查询"""
        try:
            self.queue.put(lambda: self._disable_input(True))
            # 使用text_chat方法处理响应
            response = self.bot.text_chat(query)
            self._thread_safe_display("AI", response)
            # 语音播报会在text_chat方法中自动处理
        except Exception as e:
            self._thread_safe_display("系统", f"处理错误: {str(e)}")
        finally:
            self.queue.put(lambda: self._disable_input(False))

    def _toggle_camera(self):
        """切换摄像头状态"""
        if self.bot.toggle_camera():
            self.camera_button.configure(text="关闭摄像头")
            self._display_message("系统", "摄像头已开启")
            # 启动摄像头画面更新
            self._update_camera_display()
            # 如果语音识别已开启，则开始监听
            if self.bot.is_listening:
                self.bot.start_listening(self._handle_voice_input)
        else:
            self.camera_button.configure(text="开启摄像头")
            self._display_message("系统", "摄像头已关闭")
            # 停止摄像头画面更新
            self.after_cancel(self._camera_update_id)
            # 清除摄像头显示
            self.camera_display.configure(image='')
            self.camera_display.image = None  # 清除引用
            # 如果语音识别已开启，则停止监听
            if self.bot.is_listening:
                self.bot.stop_listening()

    def _update_camera_display(self):
        """更新摄像头画面显示"""
        if self.bot.is_camera_active and self.bot.camera:
            try:
                ret, frame = self.bot.camera.read()
                if ret:
                    # 调整图像大小以适应显示区域
                    target_width = 300
                    target_height = 225  # 4:3比例
                    
                    # 确保frame不是None且有正确的形状
                    if frame is not None and len(frame.shape) == 3:
                        # 保持宽高比的调整
                        height, width = frame.shape[:2]
                        aspect_ratio = width / height
                        target_aspect_ratio = target_width / target_height
                        
                        if aspect_ratio > target_aspect_ratio:
                            # 以宽度为基准调整
                            new_width = target_width
                            new_height = int(target_width / aspect_ratio)
                        else:
                            # 以高度为基准调整
                            new_height = target_height
                            new_width = int(target_height * aspect_ratio)
                        
                        # 调整图像大小
                    frame = cv2.resize(frame, (new_width, new_height))
                    # 创建一个黑色背景
                    background = np.zeros((target_height, target_width, 3), dtype=np.uint8)
                    
                    # 计算居中位置
                    y_offset = (target_height - new_height) // 2
                    x_offset = (target_width - new_width) // 2
                    # 将调整后的图像放在背景中央
                    background[y_offset:y_offset+new_height, x_offset:x_offset+new_width] = frame
                    
                    # 转换颜色空间并显示
                    frame_rgb = cv2.cvtColor(background, cv2.COLOR_BGR2RGB)
                    image = Image.fromarray(frame_rgb)
                    photo = ImageTk.PhotoImage(image=image)
                    self.camera_display.configure(image=photo)
                    self.camera_display.image = photo  # 保持引用
                if ret:
                    if frame is not None and len(frame.shape) == 3:
                        # 前面的代码处理完图像后
                        self.camera_display.configure(image=photo)
                        self.camera_display.image = photo  # 保持引用
                    else:
                        print("无效的摄像头帧")
                else:
                    print("无法读取摄像头帧")
            except Exception as e:
                print(f"摄像头画面更新错误: {str(e)}")
                # 如果发生错误，尝试重新初始化摄像头
                self.bot.toggle_camera()
                self.bot.toggle_camera()
        
        # 继续更新
        self._camera_update_id = self.after(30, self._update_camera_display)  # 约30fps

    def _open_file_dialog(self):
        """打开文件选择对话框"""
        filetypes = (
            ('所有支持的文件', '*.pdf;*.docx;*.pptx;*.png;*.jpg;*.jpeg'),
            ('PDF文件', '*.pdf'),
            ('Word文件', '*.docx'),
            ('PPT文件', '*.pptx'),
            ('图片文件', '*.png;*.jpg;*.jpeg'),
            ('所有文件', '*.*')
        )
        
        filename = filedialog.askopenfilename(
            title='选择文件',
            filetypes=filetypes
        )
        
        if filename:
            self._process_file(filename)

    def _handle_drop(self, event):
        """处理文件拖放"""
        file_path = event.data
        # 移除可能的大括号
        file_path = file_path.strip('{}')
        self._process_file(file_path)

    def _process_file(self, file_path):
        """处理上传的文件"""
        try:
            file_name = os.path.basename(file_path)
            file_ext = os.path.splitext(file_path)[1].lower()
            
            # 更新状态为处理中
            self._update_file_status(f"正在处理文件: {file_name}")
            
            # 创建处理线程
            def process_thread():
                try:
                    if file_ext in ['.pdf']:
                        content = self._extract_pdf_content(file_path)
                    elif file_ext in ['.docx']:
                        content = self._extract_docx_content(file_path)
                    elif file_ext in ['.pptx']:
                        content = self._extract_pptx_content(file_path)
                    elif file_ext in ['.png', '.jpg', '.jpeg']:
                        content = self._process_image(file_path)
                    else:
                        raise ValueError(f"不支持的文件类型: {file_ext}")
                    
                    # 将处理结果添加到待处理列表
                    self.pending_files.append({
                        'name': file_name,
                        'content': content,
                        'type': file_ext
                    })
                    
                    # 更新状态为处理完成
                    def update_success():
                        self._update_file_status(
                        f"文件处理完成: {file_name}\n"
                        f"请在输入框输入提示词（可选）后点击发送"
                        )
                    self.queue.put(update_success)
                    
                except Exception as e:
                    # 创建一个闭包来捕获错误信息
                    def update_error(error_msg):
                        def _update():
                            self._update_file_status(error_msg)
                        return _update
                    
                    error_message = f"文件处理错误: {file_name}\n{str(e)}"
                    self.queue.put(update_error(error_message))
            
            # 启动处理线程
            threading.Thread(target=process_thread, daemon=True).start()
            
        except Exception as e:
            self._update_file_status(f"文件处理错误: {str(e)}")

    def _extract_pdf_content(self, file_path):
        """提取PDF文件内容"""
        doc = fitz.open(file_path)
        text = ""
        for page in doc:
            text += page.get_text()
        doc.close()
        return text

    def _extract_docx_content(self, file_path):
        """提取Word文件内容"""
        doc = Document(file_path)
        text = ""
        for para in doc.paragraphs:
            text += para.text + "\n"
        return text

    def _extract_pptx_content(self, file_path):
        """提取PPT文件内容"""
        prs = Presentation(file_path)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text

    def _process_image(self, file_path):
        """处理图片文件"""
        try:
            # 读取图片并转换为base64
            with Image.open(file_path) as img:
                buffered = BytesIO()
                img.save(buffered, format="PNG")
                img_str = base64.b64encode(buffered.getvalue()).decode()
            
            # 调用豆包视觉API分析图片
            return self.bot._call_doubao_vision_api(img_str, "请分析这张图片的内容")
            
        except Exception as e:
            raise Exception(f"图片处理错误: {str(e)}")

    def _update_file_status(self, message):
        """更新文件处理状态"""
        self.file_status_label.configure(text=message)

    def upload_files(self):
        """上传音频和视频文件"""
        audio_file = filedialog.askopenfilename(title="选择音频文件")
        video_file = filedialog.askopenfilename(title="选择视频文件")
        if audio_file and video_file:
            self.predict_personality(audio_file, video_file)

    def predict_personality(self, audio_file, video_file):
        """预测性格特征"""
        audio_data = self.bot.preprocess_audio(audio_file)
        video_data = self.bot.preprocess_video(video_file)
        predictions = self.bot.predict_personality(audio_data, video_data)
        self.display_predictions(predictions)

    def display_predictions(self, predictions):
        """显示预测结果"""
        result = "\n".join([f"{trait}: {score:.2f}" for trait, score in predictions.items()])
        messagebox.showinfo("预测结果", result)

    def _on_voice_button_press(self, event):
        """处理语音按钮按下事件"""
        if not self.bot.is_camera_active:
            messagebox.showwarning("警告", "进行语音识别前，请先开启摄像头。")
            return

        if not self.bot.is_listening:
            if self.bot.start_listening(self._handle_voice_input):
                self.voice_button.config(text="松开结束")
                self._thread_safe_display("系统", "正在聆听，请说话...")
            else:
                messagebox.showerror("错误", "语音识别启动失败，请检查麦克风。")

    def _on_voice_button_release(self, event):
        """处理语音按钮松开事件"""
        if self.bot.is_listening:
            self.bot.stop_listening()
            self.voice_button.config(text="按住说话")
            self._thread_safe_display("系统", "聆听结束。")

    def _start_personality_analysis(self):
        """开始性格分析"""
        try:
            # 创建临时目录
            temp_dir = os.path.join(os.path.dirname(os.path.abspath(__file__)), "temp")
            os.makedirs(temp_dir, exist_ok=True)
            
            # 使用临时目录中的完整路径
            video_path = os.path.join(temp_dir, "temp_video.mp4")
            audio_path = os.path.join(temp_dir, "temp_audio.wav")
            
            # 检查摄像头状态
            if not self.bot.is_camera_active:
                messagebox.showwarning("警告", "请先开启摄像头")
                return
            
            # 提示用户准备
            if not messagebox.askyesno("准备开始", 
                "即将开始录制10秒视频和音频用于性格分析。\n\n" + 
                "请确保：\n" +
                "1. 摄像头已开启\n" +
                "2. 麦克风已连接\n" +
                "3. 准备好对着摄像头说话\n\n" +
                "准备好了吗？"):
                return
            
            # 显示进度条和预览窗口
            self.progress_bar.pack(fill=tk.X)
            self.progress_bar["value"] = 0
            
            # 暂时关闭语音识别
            was_listening = self.bot.is_listening
            if was_listening:
                self.bot.stop_listening()
                self.voice_button.configure(text="语音识别: 关")
            
            try:
                self._show_preview_window()
            except Exception as e:
                print(f"预览窗口创建失败: {str(e)}")
                # 如果预览窗口创建失败，恢复语音识别状态
                if was_listening:
                    self.bot.start_listening(self._handle_voice_input)
                    self.voice_button.configure(text="语音识别: 开")
                return
            
            # 禁用所有控制按钮
            self._disable_input(True)
            
            # 开始录制
            self._thread_safe_display("系统", "开始录制，请自然地对着摄像头说话...\n剩余时间：10秒")
            
            # 创建录制线程
            def record_thread():
                video_writer = None
                audio_recorder = None
                
                try:
                    # 初始化视频写入器
                    fourcc = cv2.VideoWriter_fourcc(*'mp4v')
                    video_writer = cv2.VideoWriter(video_path, fourcc, 30.0, (640,480))
                    
                    # 初始化音频录制
                    import pyaudio
                    import wave
                    import numpy as np
                    
                    CHUNK = 1024
                    FORMAT = pyaudio.paFloat32
                    CHANNELS = 1
                    RATE = 16000
                    
                    p = pyaudio.PyAudio()
                    audio_recorder = p.open(format=FORMAT,
                                          channels=CHANNELS,
                                          rate=RATE,
                                          input=True,
                                          frames_per_buffer=CHUNK)
                    
                    audio_frames = []
                    start_time = time.time()
                    
                    # 录制10秒
                    while time.time() - start_time < 10:
                        try:
                            # 更新进度条
                            progress = int(((time.time() - start_time) / 10) * 100)
                            self.queue.put(lambda p=progress: self.progress_bar.configure(value=p))
                            
                            # 更新剩余时间显示
                            remaining_time = 10 - int(time.time() - start_time)
                            if remaining_time % 2 == 0:  # 每2秒更新一次显示
                                self.queue.put(lambda t=remaining_time: 
                                    self._thread_safe_display("系统", f"录制中...\n剩余时间：{t}秒"))
                            
                            # 录制视频帧
                            ret, frame = self.bot.camera.read()
                            if ret:
                                # 更新预览
                                if self.preview_window:
                                    self.queue.put(lambda f=frame.copy(): self._update_preview(f))
                                
                                # 保存视频帧
                                if video_writer is not None:
                                    video_writer.write(frame)
                            
                            # 录制音频
                            audio_data = audio_recorder.read(CHUNK, exception_on_overflow=False)
                            audio_frames.append(audio_data)
                            
                            # 更新音频预览
                            if len(audio_frames) % 10 == 0 and self.preview_window:  # 每10帧更新一次
                                audio_data_np = np.frombuffer(b''.join(audio_frames[-50:]), dtype=np.float32)
                                self.queue.put(lambda a=audio_data_np: self._update_preview(None, a))
                            
                        except Exception as frame_error:
                            print(f"帧处理错误: {str(frame_error)}")
                            continue
                    
                    # 停止录制并保存音频
                    try:
                        if audio_recorder:
                            audio_recorder.stop_stream()
                            audio_recorder.close()
                        p.terminate()
                        
                        # 将float32数据转换为int16
                        audio_data = np.frombuffer(b''.join(audio_frames), dtype=np.float32)
                        audio_data = (audio_data * 32767).astype(np.int16)
                        
                        # 确保目录存在
                        os.makedirs(os.path.dirname(audio_path), exist_ok=True)
                        
                        # 保存为WAV文件
                        with wave.open(audio_path, 'wb') as wf:
                            wf.setnchannels(CHANNELS)
                            wf.setsampwidth(2)  # 2 bytes for int16
                            wf.setframerate(RATE)
                            wf.writeframes(audio_data.tobytes())
                        
                    except Exception as audio_error:
                        self.queue.put(lambda e=audio_error: 
                            self._thread_safe_display("系统", f"音频保存失败: {str(e)}"))
                        raise
                    
                    # 显示录制完成消息
                    self.queue.put(lambda: self._thread_safe_display("系统", "录制完成！正在分析..."))
                    
                    # 确保文件已经保存
                    if not os.path.exists(audio_path) or not os.path.exists(video_path):
                        raise Exception("录制文件保存失败")
                    
                    # 关闭预览窗口
                    self.queue.put(self._close_preview)
                    
                    # 进行分析
                    try:
                        self.queue.put(lambda: self._analyze_recordings(video_path, audio_path))
                    except Exception as analysis_error:
                        self.queue.put(lambda e=analysis_error: 
                            self._thread_safe_display("系统", f"分析失败: {str(e)}"))
                        raise
                
                except Exception as e:
                    error_msg = str(e)
                    self.queue.put(lambda msg=error_msg: 
                        self._thread_safe_display("系统", f"录制失败: {msg}"))
                    # 发生错误时清理文件
                    try:
                        if os.path.exists(video_path):
                            os.remove(video_path)
                        if os.path.exists(audio_path):
                            os.remove(audio_path)
                    except Exception as cleanup_error:
                        print(f"清理临时文件失败: {str(cleanup_error)}")
                
                finally:
                    # 清理资源
                    if video_writer is not None:
                        video_writer.release()
                    
                    if audio_recorder is not None:
                        try:
                            audio_recorder.stop_stream()
                            audio_recorder.close()
                        except:
                            pass
                    
                    # 隐藏进度条
                    self.queue.put(lambda: self.progress_bar.pack_forget())
                    
                    # 恢复按钮状态
                    self.queue.put(lambda: self._disable_input(False))
                    
                    # 恢复之前的语音识别状态
                    # if was_listening: # 这部分逻辑也需要调整，因为语音识别不再是持续状态
                    #     self.queue.put(lambda: self._toggle_voice_recognition())
            
            # 启动录制线程
            threading.Thread(target=record_thread, daemon=True).start()
            
        except Exception as e:
            self._thread_safe_display("系统", f"启动录制失败: {str(e)}")
            self.progress_bar.pack_forget()
            self._disable_input(False)
            if was_listening:
                self._toggle_voice_recognition()

    def _analyze_recordings(self, video_path, audio_path):
        """分析录制的音视频"""
        try:
            # 进行性格分析
            report = self.bot.analyze_personality(audio_path, video_path)
            
            # 保存到历史记录
            self._save_to_history(report)
            
            # 创建可视化报告
            visual_report = self._create_visual_report(report)
            
            # 显示分析结果
            self._thread_safe_display("系统", "分析完成！\n\n" + visual_report)
            
            # 播放提示音
            import winsound
            winsound.Beep(1000, 500)
            
            # 弹出结果窗口
            self._show_result_window(report)
            
        except Exception as e:
            self._thread_safe_display("系统", f"分析失败: {str(e)}")
        finally:
            # 清理临时文件
            try:
                if os.path.exists(video_path):
                    os.remove(video_path)
                if os.path.exists(audio_path):
                    os.remove(audio_path)
            except Exception as cleanup_error:
                print(f"清理临时文件失败: {str(cleanup_error)}")

    def _create_visual_report(self, report_text):
        """创建可视化的报告"""
        # 解析报告文本
        lines = report_text.split('\n')
        scores = {}
        for line in lines:
            if ':' in line:
                trait, score = line.split(':')
                try:
                    scores[trait.strip()] = float(score)
                except:
                    continue
        
        # 创建可视化文本
        visual = "性格特征分析报告：\n\n"
        for trait, score in scores.items():
            bar_length = int(score * 20)  # 将0-1的分数转换为0-20的长度
            bar = '█' * bar_length + '░' * (20 - bar_length)
            visual += f"{trait}:\n{bar} {score:.2f}\n\n"
        
        return visual

    def _save_to_history(self, report):
        """保存分析结果到历史记录"""
        try:
            # 创建历史记录目录
            os.makedirs("history", exist_ok=True)
            
            # 生成文件名
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            filename = f"history/personality_analysis_{timestamp}.txt"
            
            # 保存报告
            with open(filename, "w", encoding="utf-8") as f:
                f.write(f"分析时间: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}\n\n")
                f.write(report)
                
        except Exception as e:
            print(f"保存历史记录失败: {str(e)}")

    def _show_history(self):
        """显示历史记录"""
        try:
            # 检查历史记录目录
            if not os.path.exists("history"):
                messagebox.showinfo("提示", "暂无历史记录")
                return
            
            # 创建历史记录窗口
            history_window = tk.Toplevel(self)
            history_window.title("历史记录")
            history_window.geometry("600x400")
            
            # 创建列表框
            listbox = tk.Listbox(history_window, font=('Microsoft YaHei', 10))
            listbox.pack(fill=tk.BOTH, expand=True, padx=5, pady=5)
            
            # 添加滚动条
            scrollbar = ttk.Scrollbar(listbox)
            scrollbar.pack(side=tk.RIGHT, fill=tk.Y)
            listbox.config(yscrollcommand=scrollbar.set)
            scrollbar.config(command=listbox.yview)
            
            # 读取历史记录
            history_files = sorted(os.listdir("history"), reverse=True)
            for file in history_files:
                if file.startswith("personality_analysis_"):
                    with open(f"history/{file}", "r", encoding="utf-8") as f:
                        content = f.read()
                    listbox.insert(tk.END, content)
                    listbox.insert(tk.END, "\n" + "="*50 + "\n")
            
        except Exception as e:
            messagebox.showerror("错误", f"读取历史记录失败: {str(e)}")

    def _show_preview_window(self):
        """显示录制预览窗口"""
        if self.preview_window is None:
            self.preview_window = tk.Toplevel(self)
            self.preview_window.title("录制预览")
            self.preview_window.geometry("800x600")
            
            # 创建预览画布
            preview_frame = ttk.Frame(self.preview_window)
            preview_frame.pack(fill=tk.BOTH, expand=True)
            
            # 创建视频预览区域
            video_frame = ttk.LabelFrame(preview_frame, text="视频预览")
            video_frame.pack(side=tk.TOP, fill=tk.BOTH, expand=True, padx=5, pady=5)
            
            self.preview_canvas = tk.Canvas(
                video_frame,
                width=640,
                height=480,
                bg='black'
            )
            self.preview_canvas.pack(padx=5, pady=5)
            
            # 添加音频波形图
            audio_frame = ttk.LabelFrame(preview_frame, text="音频波形")
            audio_frame.pack(side=tk.TOP, fill=tk.X, padx=5, pady=5)
            
            self.preview_fig = plt.Figure(figsize=(8, 2))
            self.preview_fig.patch.set_facecolor('#f0f0f0')
            
            # 设置中文字体
            plt.rcParams['font.sans-serif'] = ['SimHei']
            plt.rcParams['axes.unicode_minus'] = False
            
            canvas = FigureCanvasTkAgg(self.preview_fig, audio_frame)
            canvas.get_tk_widget().pack(fill=tk.X, padx=5, pady=5)
            
            # 添加控制按钮框架
            control_frame = ttk.Frame(preview_frame)
            control_frame.pack(side=tk.BOTTOM, fill=tk.X, padx=5, pady=5)
            
            # 添加关闭按钮
            close_button = ttk.Button(
                control_frame,
                text="关闭预览",
                command=self._close_preview
            )
            close_button.pack(side=tk.RIGHT, padx=5)
            
            # 添加帧率显示
            self.fps_label = ttk.Label(control_frame, text="FPS: 0")
            self.fps_label.pack(side=tk.LEFT, padx=5)
            
            # 初始化FPS计算
            self.frame_times = []
            self.last_fps_update = time.time()
            
            self.preview_window.protocol("WM_DELETE_WINDOW", self._close_preview)
            
    def _update_preview(self, frame, audio_data=None):
        """更新预览窗口的内容"""
        if self.preview_window and self.preview_canvas:
            try:
                # 更新视频预览
                if frame is not None:
                    # 计算FPS
                    current_time = time.time()
                    self.frame_times.append(current_time)
                    
                    # 只保留最近1秒的帧时间
                    while self.frame_times and self.frame_times[0] < current_time - 1:
                        self.frame_times.pop(0)
                    
                    # 每秒更新一次FPS显示
                    if current_time - self.last_fps_update >= 1:
                        fps = len(self.frame_times)
                        self.fps_label.config(text=f"FPS: {fps}")
                        self.last_fps_update = current_time
                    
                    # 调整视频帧大小
                    frame = cv2.resize(frame, (640, 480))
                    frame = cv2.cvtColor(frame, cv2.COLOR_BGR2RGB)
                    
                    # 使用PIL进行图像处理
                    img = Image.fromarray(frame)
                    img = ImageTk.PhotoImage(image=img)
                    
                    # 更新画布
                    self.preview_canvas.create_image(
                        self.preview_canvas.winfo_width() // 2,
                        self.preview_canvas.winfo_height() // 2,
                        image=img,
                        anchor=tk.CENTER
                    )
                    self.preview_canvas.image = img
                
                # 更新音频波形图
                if audio_data is not None and self.preview_fig is not None:
                    self.preview_fig.clear()
                    ax = self.preview_fig.add_subplot(111)
                    
                    # 计算音频包络
                    window_size = 100
                    envelope = np.abs(audio_data)
                    envelope = np.convolve(envelope, np.ones(window_size)/window_size, mode='same')
                    
                    # 绘制波形
                    ax.plot(envelope, color='#2196F3', linewidth=1)
                    ax.fill_between(range(len(envelope)), envelope, color='#2196F3', alpha=0.2)
                    
                    # 设置样式
                    ax.set_title("音频波形", fontproperties='SimHei', pad=10)
                    ax.set_xticks([])
                    ax.set_yticks([])
                    ax.spines['top'].set_visible(False)
                    ax.spines['right'].set_visible(False)
                    ax.spines['left'].set_visible(False)
                    
                    # 更新画布
                    self.preview_fig.canvas.draw()
                    
            except Exception as e:
                print(f"更新预览失败: {str(e)}")
                
    def _close_preview(self):
        """关闭预览窗口"""
        if self.preview_window:
            self.recording_preview = False
            self.preview_window.destroy()
            self.preview_window = None
            self.preview_canvas = None
            self.preview_fig = None
            self.frame_times = []

    def _show_result_window(self, report):
        """显示分析结果窗口"""
        result_window = tk.Toplevel(self)
        result_window.title("性格分析报告")
        result_window.geometry("800x600")
        
        # 创建notebook用于多标签页显示
        notebook = ttk.Notebook(result_window)
        notebook.pack(fill=tk.BOTH, expand=True, padx=5, pady=5)
        
        # 概述标签页
        summary_frame = ttk.Frame(notebook)
        notebook.add(summary_frame, text="概述")
        
        # 解析报告文本
        scores = {}
        trait_descriptions = {
            "开放性": "反映个人对新体验的接受程度，包括创造力、好奇心和艺术兴趣。",
            "尽责性": "反映个人的组织能力、可靠性和责任心。",
            "外向性": "反映个人在社交场合的表现和能量水平。",
            "亲和性": "反映个人与他人相处的态度和同理心。",
            "神经质": "反映个人的情绪稳定性和压力应对能力。"
        }
        
        # MBTI结果和总结
        mbti_result = ""
        mbti_desc = ""
        model_analysis = ""
        
        # 分段解析报告
        lines = report.split('\n')
        in_deepseek_section = False
        
        for i, line in enumerate(lines):
            # 提取五大人格得分
            if ':' in line:
                parts = line.split(':', 1)  # 只分割第一个冒号
                if len(parts) == 2:
                    trait, score = parts
                    trait_name = trait.strip()
                    # 提取MBTI结果
                    if "MBTI性格类型预测" in trait_name:
                        mbti_result = score.strip()
                    # 提取MBTI描述
                    elif "类型说明" in trait_name:
                        mbti_desc = score.strip()
                    # 只提取五大分析结果
                    elif any(key in trait_name for key in ["外向性", "神经质", "尽责性", "宜人性", "开放性"]):
                        try:
                            scores[trait_name] = float(score)
                        except ValueError:
                            pass
            
            # 提取大模型分析结果
            if "大模型分析结果：" in line:
                in_deepseek_section = True
                continue
            
            # 累积大模型分析内容
            if in_deepseek_section and line.strip() and not line.startswith("获取深度分析时发生错误"):
                model_analysis += line + "\n"
        
        # 创建概述页面内容
        summary_content = ttk.Frame(summary_frame)
        summary_content.pack(fill=tk.BOTH, expand=True, padx=20, pady=20)
        
        # 添加标题
        title_label = ttk.Label(
            summary_content, 
            text="性格分析报告概述", 
            font=("Microsoft YaHei", 16, "bold")
        )
        title_label.pack(pady=(0, 20))
        
        # 添加五大人格得分概要
        if scores:
            personality_frame = ttk.LabelFrame(summary_content, text="五大人格特质")
            personality_frame.pack(fill=tk.X, pady=10)
            
            # 为每个特质创建进度条显示
            for trait, score in scores.items():
                trait_frame = ttk.Frame(personality_frame)
                trait_frame.pack(fill=tk.X, pady=5, padx=10)
                
                # 特质名称
                trait_label = ttk.Label(trait_frame, text=f"{trait}:", width=20)
                trait_label.pack(side=tk.LEFT, padx=(0, 10))
                
                # 进度条
                pb = ttk.Progressbar(
                    trait_frame, 
                    orient=tk.HORIZONTAL, 
                    length=300, 
                    mode='determinate', 
                    value=score*100
                )
                pb.pack(side=tk.LEFT, padx=5)
                
                # 分数值
                score_label = ttk.Label(trait_frame, text=f"{score:.2f}")
                score_label.pack(side=tk.LEFT, padx=5)
        
        # 添加MBTI结果
        if mbti_result:
            mbti_frame = ttk.LabelFrame(summary_content, text="MBTI人格类型")
            mbti_frame.pack(fill=tk.X, pady=10)
            
            mbti_content = ttk.Frame(mbti_frame)
            mbti_content.pack(fill=tk.X, padx=10, pady=10)
            
            mbti_label = ttk.Label(
                mbti_content, 
                text=mbti_result, 
                font=("Microsoft YaHei", 12, "bold")
            )
            mbti_label.pack(side=tk.LEFT, padx=(0, 10))
            
            if mbti_desc:
                mbti_desc_label = ttk.Label(
                    mbti_content, 
                    text=f"- {mbti_desc}",
                    wraplength=500
                )
                mbti_desc_label.pack(side=tk.LEFT)
        
        # 添加大模型分析
        if model_analysis.strip():
            analysis_frame = ttk.LabelFrame(summary_content, text="AI深度分析")
            analysis_frame.pack(fill=tk.BOTH, expand=True, pady=10)
            
            analysis_text = scrolledtext.ScrolledText(
                analysis_frame, 
                wrap=tk.WORD,
                height=10
            )
            analysis_text.pack(fill=tk.BOTH, expand=True, padx=10, pady=10)
            analysis_text.insert(tk.END, model_analysis)
            analysis_text.config(state=tk.DISABLED)
        
        # 雷达图标签页
        radar_frame = ttk.Frame(notebook)
        notebook.add(radar_frame, text="性格雷达图")
        
        if scores:
            # 创建一个matplotlib figure
            fig = plt.Figure(figsize=(5, 4), dpi=100)
            ax = fig.add_subplot(111, polar=True)
            
            # 准备数据
            traits = list(scores.keys())
            values = list(scores.values())
            num_traits = len(traits)
            
            if num_traits > 0:
                # 设置雷达图的角度
                angles = np.linspace(0, 2*np.pi, num_traits, endpoint=False).tolist()
                
                # 闭合雷达图
                values += [values[0]]
                angles += [angles[0]]
                
                # 绘制雷达图
                ax.plot(angles, values, 'o-', linewidth=2)
                ax.fill(angles, values, alpha=0.25)
                
                # 设置刻度标签
                ax.set_xticks(angles[:-1])
                ax.set_xticklabels(traits)
                
                # 设置y轴范围为0到1
                ax.set_ylim(0, 1)
                
                # 添加网格
                ax.grid(True)
                
                # 设置标题
                ax.set_title("性格五大维度分析")
                
                # 将图形嵌入到tkinter中
                canvas = FigureCanvasTkAgg(fig, master=radar_frame)
                canvas.draw()
                canvas.get_tk_widget().pack(side=tk.TOP, fill=tk.BOTH, expand=1)
                
                # 添加matplotlib工具栏
                toolbar = NavigationToolbar2Tk(canvas, radar_frame)
                toolbar.update()
                canvas.get_tk_widget().pack(side=tk.TOP, fill=tk.BOTH, expand=1)
        
        # 详细报告标签页
        report_frame = ttk.Frame(notebook)
        notebook.add(report_frame, text="详细报告")
        
        # 创建文本区域显示完整报告
        report_text = scrolledtext.ScrolledText(report_frame, wrap=tk.WORD)
        report_text.pack(fill=tk.BOTH, expand=True, padx=5, pady=5)
        report_text.insert(tk.END, report)
        report_text.config(state=tk.DISABLED)
        
        # 添加保存按钮
        save_frame = ttk.Frame(result_window)
        save_frame.pack(fill=tk.X, padx=5, pady=5)
        
        save_button = ttk.Button(
            save_frame,
            text="保存报告",
            command=lambda: self._save_report(report, scores)
        )
        save_button.pack(side=tk.RIGHT, padx=5)

    def _get_trait_analysis(self, trait, score):
        """获取特定性格特征的详细分析"""
        if trait == "开放性":
            if score > 0.7:
                return "您表现出很高的开放性，喜欢尝试新事物，富有创造力和想象力。"
            elif score > 0.4:
                return "您的开放性处于中等水平，能够在保持稳定的同时接受适度的改变。"
            else:
                return "您倾向于保守和传统，更喜欢熟悉的环境和经验。"
        elif trait == "尽责性":
            if score > 0.7:
                return "您非常负责任，有很强的计划性和组织能力。"
            elif score > 0.4:
                return "您的责任心处于中等水平，能够在需要时保持专注和努力。"
            else:
                return "您可能需要提高自己的计划性和组织能力。"
        # ... 其他特征的分析 ...
        return ""

    def _get_trait_suggestions(self, trait, score):
        """获取特定性格特征的改进建议"""
        if trait == "开放性":
            if score < 0.4:
                return "• 尝试接触新的活动和体验\n• 培养创造性爱好\n• 学习新的技能"
            elif score < 0.7:
                return "• 在保持现有兴趣的同时，适度尝试新事物\n• 参与创造性活动"
            else:
                return "• 保持好奇心和创造力\n• 注意在探索新事物时保持适度"
        elif trait == "尽责性":
            if score < 0.4:
                return "• 制定日常计划和目标\n• 培养时间管理能力\n• 建立工作检查清单"
            elif score < 0.7:
                return "• 继续保持良好的工作习惯\n• 适当提高计划的执行力"
            else:
                return "• 保持良好的组织能力\n• 注意在追求完美时保持弹性"
        # ... 其他特征的建议 ...
        return ""

    def _save_report(self, report, scores):
        """保存分析报告"""
        try:
            # 选择保存位置
            file_path = filedialog.asksaveasfilename(
                defaultextension=".pdf",
                filetypes=[("PDF文件", "*.pdf"), ("文本文件", "*.txt")],
                title="保存报告"
            )
            
            if file_path:
                if file_path.endswith('.pdf'):
                    self._save_pdf_report(file_path, report, scores)
                else:
                    self._save_text_report(file_path, report)
                    
                messagebox.showinfo("成功", "报告已保存")
                
        except Exception as e:
            messagebox.showerror("错误", f"保存报告失败: {str(e)}")

    def _save_pdf_report(self, file_path, report, scores):
        """保存PDF格式的报告"""
        from reportlab.lib import colors
        from reportlab.lib.pagesizes import A4
        from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
        from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
        from reportlab.pdfbase import pdfmetrics
        from reportlab.pdfbase.ttfonts import TTFont
        
        # 注册中文字体
        pdfmetrics.registerFont(TTFont('SimSun', 'simsun.ttc'))
        
        # 创建PDF文档
        doc = SimpleDocTemplate(file_path, pagesize=A4)
        styles = getSampleStyleSheet()
        
        # 创建中文样式
        styles.add(ParagraphStyle(
            name='Chinese',
            fontName='SimSun',
            fontSize=12,
            leading=14
        ))
        
        # 准备内容
        content = []
        
        # 添加标题
        content.append(Paragraph("性格分析报告", styles['Title']))
        content.append(Spacer(1, 12))
        
        # 添加报告内容
        for line in report.split('\n'):
            if line.strip():
                content.append(Paragraph(line, styles['Chinese']))
                content.append(Spacer(1, 6))
        
        # 生成PDF
        doc.build(content)

    def _save_text_report(self, file_path, report):
        """保存文本格式的报告"""
        with open(file_path, 'w', encoding='utf-8') as f:
            f.write(report)

    def _toggle_tts(self):
        """切换语音播报开关"""
        if hasattr(self.bot, 'toggle_tts'): 
            is_enabled = self.bot.toggle_tts()
            self.tts_button.config(text=f"语音播报: {'开' if is_enabled else '关'}")
        else:
            messagebox.showwarning("功能缺失", "TTS切换功能未在机器人中实现。")

    def _thread_safe_display(self, sender, message):
        """线程安全的消息显示方法"""
        if hasattr(self, 'queue') and self.queue is not None:
             self.queue.put(lambda: self._display_message(sender, message))
        else: 
             self._display_message(sender, message)

    def _display_message(self, sender, message):
        """实际执行消息显示"""
        if not hasattr(self, 'chat_area') or self.chat_area is None:
            print(f"Debug: Chat area not initialized. Msg: {sender} - {message}")
            return

        self.chat_area.config(state='normal')
        timestamp = datetime.now().strftime("%H:%M:%S")
        
        sender_tag_name = f"sender_{sender.replace(' ', '')}" # Create a unique tag name for sender

        if sender == "系统":
            color = "#666666"
        elif sender == "您":
            color = "#2196F3"
        else: # AI
            color = "#4CAF50"
            
        self.chat_area.insert(tk.END, f"\\n[{timestamp}] ", "timestamp")
        self.chat_area.tag_configure("timestamp", foreground="#999999")
        
        self.chat_area.insert(tk.END, f"{sender}：\\n", sender_tag_name) 
        self.chat_area.tag_configure(sender_tag_name, foreground=color, font=('Microsoft YaHei', 10, 'bold'))
        
        self.chat_area.insert(tk.END, f"{message}\\n", "message_tag") # General tag for message content
        self.chat_area.tag_configure("message_tag", foreground="#333333")
        
        self.chat_area.see(tk.END)
        self.chat_area.config(state='disabled')

    def _disable_input(self, disabled):
        """在查询处理期间禁用/启用文本输入和发送按钮。"""
        text_input_state = 'disabled' if disabled else 'normal'
        send_button_state = 'disabled' if disabled else 'normal'

        if hasattr(self, 'input_text'):
            self.input_text.config(state=text_input_state)
        
        if hasattr(self, 'send_button'):
            self.send_button.config(state=send_button_state)
        
        # "按住说话" 按钮也应在处理查询时禁用
        # 并在处理完成后根据登录状态恢复
        if hasattr(self, 'voice_button'):
            if disabled:
                self.voice_button.config(state='disabled')
            else:
                # 恢复状态时，如果用户已登录，则设为normal，否则保持disabled (或根据未登录时的默认状态)
                # 这里的逻辑应该与 _update_ui_for_login_status 对 voice_button 的设置一致
                current_voice_button_state = 'normal' if self.logged_in_user else 'disabled' 
                # However, _update_ui_for_login_status might be more comprehensive for this.
                # For now, let's simplify and assume if input is enabled, voice might be too if logged in.
                if self.logged_in_user: # Only enable if user is logged in and input is being re-enabled
                    self.voice_button.config(state='normal')
                else: # If not logged in, or if it should be specifically disabled even when input is enabled
                    self.voice_button.config(state='disabled') # Or consult _update_ui_for_login_status logic